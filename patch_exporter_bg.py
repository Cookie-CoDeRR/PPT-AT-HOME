import re

with open('backend/services/export_pptx.py', 'r') as f:
    code = f.read()

imports = """
import re
from PIL import Image, ImageDraw
import base64
import tempfile
import os
"""

if "from PIL import Image" not in code:
    code = code.replace("import base64", imports)

bg_logic = """
def apply_custom_background(slide, bg_data, prs):
    bg_type = bg_data.get('type')
    bg_value = bg_data.get('value', '')
    
    # 1. Background Fill
    if bg_type == 'solid':
        color_rgb = hex_to_rgb(bg_value)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color_rgb
    
    elif bg_type == 'image' and bg_value.startswith('data:image'):
        try:
            head, data = bg_value.split(',', 1)
            img_data = base64.b64decode(data)
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            temp_file.write(img_data)
            temp_file.close()
            
            # Add full-screen picture as background (send to back)
            pic = slide.shapes.add_picture(temp_file.name, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            os.remove(temp_file.name)
            
            # Move to back (z-order 0)
            pic._element.addprevious(pic._element.getparent()[0])
            pic._element.getparent().insert(0, pic._element)
        except Exception as e:
            print("Failed to set image background:", e)
            
    elif bg_type == 'gradient' and bg_value:
        try:
            # Simple gradient generator using Pillow
            # e.g., 'linear-gradient(135deg, #0f172a 0%, #1e1b4b 100%)'
            colors = re.findall(r'#([A-Fa-f0-9]{6})', bg_value)
            if len(colors) >= 2:
                c1 = hex_to_rgb(colors[0])
                c2 = hex_to_rgb(colors[1])
                
                width = 1280
                height = 720
                img = Image.new('RGB', (width, height))
                draw = ImageDraw.Draw(img)
                
                # Draw vertical gradient approximation for simplicity
                for y in range(height):
                    r = int(c1[0] + (c2[0] - c1[0]) * y / height)
                    g = int(c1[1] + (c2[1] - c1[1]) * y / height)
                    b = int(c1[2] + (c2[2] - c1[2]) * y / height)
                    draw.line([(0, y), (width, y)], fill=(r, g, b))
                
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                img.save(temp_file.name)
                temp_file.close()
                
                pic = slide.shapes.add_picture(temp_file.name, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                os.remove(temp_file.name)
                
                pic._element.addprevious(pic._element.getparent()[0])
                pic._element.getparent().insert(0, pic._element)
        except Exception as e:
            print("Failed to set gradient background:", e)

    # 2. Overlay Logic
    if bg_type in ['image', 'gradient'] and bg_data.get('overlayOpacity', 0) > 0:
        overlay_color = bg_data.get('overlayColor', '#000000')
        opacity = bg_data.get('overlayOpacity', 0.5)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        set_shape_color(rect, overlay_color.replace('#', ''))
        rect.line.fill.background()
        
        # In PowerPoint, transparency is set on the fill format
        # However, python-pptx doesn't expose transparency directly on solid fill without XML manipulation.
        # We will use an XML hack to set alpha on the solid fill.
        try:
            fill = rect.fill
            fill.solid()
            # Alpha is 0-100000 (100% opaque = 100000, 0% = 0)
            alpha_val = int((1.0 - opacity) * 100000)
            
            # The XML structure for solid fill with alpha:
            # <a:solidFill> <a:srgbClr val="RRGGBB"> <a:alpha val="50000"/> </a:srgbClr> </a:solidFill>
            from pptx.oxml import parse_xml
            alpha_xml = f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="{alpha_val}"/>'
            
            srgbClr = rect.fill._xPr.solidFill.srgbClr
            if srgbClr is not None:
                srgbClr.append(parse_xml(alpha_xml))
        except Exception as e:
            print("Transparency XML injection failed:", e)

        # Move to back, but in front of background
        rect._element.addprevious(rect._element.getparent()[1])
        rect._element.getparent().insert(1, rect._element)
"""

if "def apply_custom_background" not in code:
    code = code.replace("def export_presentation", bg_logic + "\n\ndef export_presentation")

# Add args.custom_bg
arg_parser_addition = """
    parser.add_argument("--custom_bg", required=False, help="Custom Background JSON string")
"""
if "--custom_bg" not in code:
    code = code.replace('parser.add_argument("--output", required=True, help="Output PPTX file path")', 
                        'parser.add_argument("--output", required=True, help="Output PPTX file path")\n' + arg_parser_addition)

# Modify export_presentation signature and call
export_sig = 'def export_presentation(data, output_path, custom_bg=None):'
code = re.sub(r'def export_presentation\(data, output_path\):', export_sig, code)

main_call = """
    custom_bg_obj = None
    if args.custom_bg:
        custom_bg_obj = json.loads(args.custom_bg)
        
    export_presentation(data, args.output, custom_bg_obj)
"""
code = re.sub(r'export_presentation\(data, args.output\)', main_call.strip(), code)

# Apply it in the loop
apply_call = """
        if custom_bg:
            apply_custom_background(slide, custom_bg, prs)
"""
code = code.replace('slide = prs.slides.add_slide(blank_slide_layout)', 
                    'slide = prs.slides.add_slide(blank_slide_layout)' + apply_call)

with open('backend/services/export_pptx.py', 'w') as f:
    f.write(code)
print("Updated export_pptx.py")
