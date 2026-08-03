import re

with open('backend/services/export_pptx.py', 'r') as f:
    code = f.read()

# Add imports
imports = """from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
"""
if "CategoryChartData" not in code:
    code = code.replace("from pptx.enum.shapes import MSO_SHAPE\n", "from pptx.enum.shapes import MSO_SHAPE\n" + imports)

# Add new builders before build_default_slide
new_builders = """
def build_chart_slide(slide, data, theme, img_path, chart_type):
    add_text_box(slide, data.get("title", ""), 0.5, 0.5, 5.0, 1.0, 32, theme['titleColor'], PP_ALIGN.LEFT, True, theme.get('fontFace'))
    add_text_box(slide, data.get("description", ""), 0.5, 1.5, 5.0, 2.0, 18, theme['textColor'], PP_ALIGN.LEFT, False, theme.get('bodyFontFace'))
    
    chart_data_obj = data.get("chart_data", {})
    labels = chart_data_obj.get("labels", [])
    values = chart_data_obj.get("values", [])
    
    if not labels or not values:
        # Fallback to default
        build_default_slide(slide, data, theme, img_path)
        return
        
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series('Series 1', values)
    
    x, y, cx, cy = Inches(6.0), Inches(1.5), Inches(6.5), Inches(5.0)
    
    ctype = XL_CHART_TYPE.PIE if chart_type == 'pie' else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

def build_data_table_slide(slide, data, theme, img_path):
    add_text_box(slide, data.get("title", ""), 0.5, 0.4, 12.33, 1, 32, theme['titleColor'], PP_ALIGN.CENTER, True, theme.get('fontFace'))
    
    table_data = data.get("table_data", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    
    if not headers or not rows:
        build_default_slide(slide, data, theme, img_path)
        return
        
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    left, top, width, height = Inches(1.0), Inches(1.5), Inches(11.33), Inches(5.0)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(theme['accent'])
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = hex_to_rgb(theme['bkgd'])
            if theme.get('fontFace'):
                p.font.name = theme.get('fontFace')
            
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            if c_idx < num_cols:
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                cell.fill.solid()
                if r_idx % 2 == 0:
                    cell.fill.fore_color.rgb = hex_to_rgb(theme['shapeFill'])
                else:
                    cell.fill.fore_color.rgb = hex_to_rgb(theme['bkgd'])
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = hex_to_rgb(theme['textColor'])
                    if theme.get('bodyFontFace'):
                        p.font.name = theme.get('bodyFontFace')

def build_bento_grid_slide(slide, data, theme, img_path):
    add_text_box(slide, data.get("title", ""), 0.5, 0.4, 12.33, 1, 34, theme['titleColor'], PP_ALIGN.CENTER, True, theme.get('fontFace'))
    
    items = data.get("items", [])
    if not items:
        build_default_slide(slide, data, theme, img_path)
        return
        
    # Fixed 3-column grid mapping
    cols = 3
    col_w = (12.33 - 1.0 - (0.4 * (cols - 1))) / cols
    row_h = 2.0
    
    current_x = 0.5
    current_y = 1.5
    
    for item in items:
        size = item.get('size', 'small')
        
        # Calculate width and height based on size
        w = col_w
        h = row_h
        
        if size == 'large':
            w = (col_w * 2) + 0.4
            h = (row_h * 2) + 0.4
        elif size == 'wide':
            w = (col_w * 3) + 0.8
        elif size == 'tall':
            h = (row_h * 2) + 0.4
            
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(current_x), Inches(current_y), Inches(w), Inches(h))
        set_shape_color(rect, theme['shapeFill'])
        
        add_text_box(slide, item.get("title", ""), current_x + 0.2, current_y + 0.2, w - 0.4, 0.8, 20 if size == 'large' else 16, theme['titleColor'], PP_ALIGN.LEFT, True, theme.get('fontFace'))
        add_text_box(slide, item.get("desc", ""), current_x + 0.2, current_y + 1.0, w - 0.4, h - 1.2, 14 if size == 'large' else 12, theme['textColor'], PP_ALIGN.LEFT, False, theme.get('bodyFontFace'))
        
        # Simple flow layout (not true masonry for simplicity)
        current_x += w + 0.4
        if current_x > 12.0:
            current_x = 0.5
            current_y += h + 0.4

def build_metric_dashboard_slide(slide, data, theme, img_path):
    add_text_box(slide, data.get("title", ""), 0.5, 0.4, 12.33, 1, 34, theme['titleColor'], PP_ALIGN.CENTER, True, theme.get('fontFace'))
    
    metrics = data.get("metrics", [])
    if not metrics:
        build_default_slide(slide, data, theme, img_path)
        return
        
    box_w = 5.5
    box_h = 2.5
    
    for i, m in enumerate(metrics[:4]):
        r = i // 2
        c = i % 2
        x = 0.8 + (c * (box_w + 0.5))
        y = 1.5 + (r * (box_h + 0.5))
        
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
        set_shape_color(rect, theme['shapeFill'])
        
        add_text_box(slide, m.get("change", ""), x + box_w - 1.5, y + 0.2, 1.3, 0.5, 14, theme['accent'], PP_ALIGN.RIGHT, True, theme.get('fontFace'))
        add_text_box(slide, m.get("value", ""), x + 0.2, y + 0.7, box_w - 0.4, 1.0, 48, theme['titleColor'], PP_ALIGN.CENTER, True, theme.get('fontFace'))
        add_text_box(slide, m.get("label", ""), x + 0.2, y + 1.8, box_w - 0.4, 0.5, 16, theme['textColor'], PP_ALIGN.CENTER, True, theme.get('fontFace'))

"""
if "build_chart_slide" not in code:
    code = code.replace("def build_default_slide", new_builders + "\ndef build_default_slide")

dispatcher = """        if stype == "comparison":
            build_comparison_slide(slide, sdata, theme, img_path)
        elif stype == "timeline":
            build_timeline_slide(slide, sdata, theme, img_path)
        elif stype == "stat_callout":
            build_stat_slide(slide, sdata, theme, img_path)
        elif stype == "grid_list":
            build_grid_slide(slide, sdata, theme, img_path)
        elif stype == "chart_pie":
            build_chart_slide(slide, sdata, theme, img_path, "pie")
        elif stype == "chart_bar":
            build_chart_slide(slide, sdata, theme, img_path, "bar")
        elif stype == "data_table":
            build_data_table_slide(slide, sdata, theme, img_path)
        elif stype == "bento_grid":
            build_bento_grid_slide(slide, sdata, theme, img_path)
        elif stype == "metric_dashboard":
            build_metric_dashboard_slide(slide, sdata, theme, img_path)
        else:"""

code = re.sub(r'        if stype == "comparison":[\s\S]*?        else:', dispatcher, code)

with open('backend/services/export_pptx.py', 'w') as f:
    f.write(code)

print("Patched export_pptx.py")
