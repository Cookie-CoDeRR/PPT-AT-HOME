import json
import argparse
import os
import re
import base64
import math
import tempfile
from PIL import Image, ImageDraw
from pptx import Presentation
from density_optimizer import optimize_presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ─────────────────────────────────────────────────────────────────────────────
# CANVAS CONSTANTS  (13.333" × 7.5" = standard 16:9)
# ─────────────────────────────────────────────────────────────────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN  = 0.5          # outer safe-zone margin
CONTENT_X = MARGIN
CONTENT_W = SLIDE_W - MARGIN * 2   # 12.333"
TITLE_Y   = 0.3
TITLE_H   = 0.75
BODY_Y    = TITLE_Y + TITLE_H + 0.1   # 1.15"
BODY_H    = SLIDE_H - BODY_Y - MARGIN  # ~5.85"

def calc_title_layout(title_text, base_font_pt=24):
    """Dynamically compute title box height and body top position to prevent overlapping."""
    words = len(str(title_text).split())
    lines = math.ceil(words / 9) if words > 0 else 1
    h = max(0.6, lines * (base_font_pt / 72.0 * 1.3))
    body_y = TITLE_Y + h + 0.1
    body_h = SLIDE_H - body_y - MARGIN
    return h, body_y, body_h



# ─────────────────────────────────────────────────────────────────────────────
# COLOUR UTILITIES
# ─────────────────────────────────────────────────────────────────────────────
def hex_to_rgb(hex_str):
    hex_str = str(hex_str).lstrip('#').strip()
    if len(hex_str) != 6:
        return RGBColor(0, 0, 0)
    try:
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    except Exception:
        return RGBColor(0, 0, 0)


def theme_color(theme, key, fallback):
    """Safely read a colour from the theme dict."""
    raw = theme.get(key, fallback)
    return str(raw).lstrip('#')


def set_shape_fill(shape, hex_str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_str)


def set_shape_no_line(shape):
    try:
        shape.line.fill.background()
    except Exception:
        pass


# ─────────────────────────────────────────────────────────────────────────────
# PLACEHOLDER PURGE — Remove inherited master placeholders that ghost-render
# ─────────────────────────────────────────────────────────────────────────────
def purge_placeholders(slide):
    """
    After adding a blank slide, python-pptx may still inherit placeholder
    shapes from the slide master (e.g. title, subtitle boxes).  These ghost
    elements render as empty boxes that show through the PPTX even though
    the slide layout is 'blank'.  This function strips them out.
    """
    sp_tree = slide.shapes._spTree
    # Collect all <p:sp> nodes that are placeholders
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    to_remove = [
        sp for sp in sp_tree.findall(f'{{{ns}}}sp')
        if sp.find(f'.//{{{ns}}}ph') is not None
    ]
    for sp in to_remove:
        sp_tree.remove(sp)
    if to_remove:
        print(f'[Purge] Removed {len(to_remove)} inherited placeholder(s) from slide.')


def add_text_box(slide, text, left, top, width, height,
                 font_size, color_hex,
                 align=PP_ALIGN.LEFT, bold=False, font_name="Calibri"):
    """Add a single-paragraph text box. All coords in inches."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.bold = bold
    p.font.size = Pt(font_size)
    p.font.color.rgb = hex_to_rgb(color_hex)
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullets(slide, bullets, left, top, width, height,
                font_size, color_hex,
                prefix="", font_name="Calibri", align=PP_ALIGN.LEFT,
                line_spacing_pt=None):
    """Add a multi-paragraph text box for a bullet list."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{prefix}{bullet}" if prefix else str(bullet)
        p.font.size = Pt(font_size)
        p.font.color.rgb = hex_to_rgb(color_hex)
        p.font.name = font_name
        p.alignment = align
        p.level = 0
        if line_spacing_pt:
            from pptx.util import Pt as _Pt
            p.line_spacing = _Pt(line_spacing_pt)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_hex, line_hex=None, line_pt=1.5):
    """Add a rounded rectangle card."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    set_shape_fill(rect, fill_hex)
    if line_hex:
        rect.line.color.rgb = hex_to_rgb(line_hex)
        rect.line.width = Pt(line_pt)
    else:
        set_shape_no_line(rect)
    return rect


def add_rectangle(slide, left, top, width, height, fill_hex):
    """Add a plain rectangle (no rounded corners)."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    set_shape_fill(rect, fill_hex)
    set_shape_no_line(rect)
    return rect


def add_oval(slide, left, top, diameter, fill_hex):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    set_shape_fill(oval, fill_hex)
    set_shape_no_line(oval)
    return oval


# ─────────────────────────────────────────────────────────────────────────────
# DEFENSIVE FALLBACKS
# ─────────────────────────────────────────────────────────────────────────────
def apply_defensive_fallbacks(slides):
    for slide in slides:
        if not slide.get("title") or not str(slide.get("title")).strip():
            slide["title"] = "Slide Content Missing"
        stype = slide.get("slide_type", "default")
        # Issue 4 fix: title_hero MUST NOT get bullet fallbacks — it only renders title+subtitle
        if stype in ["standard_text", "summary_takeaways", "stat_callout",
                     "two_column_image", "default"]:
            if not slide.get("bullets") or len(slide.get("bullets", [])) == 0:
                slide["bullets"] = ["Content will appear here", "Generated by the AI engine"]
        elif stype == "comparison":
            for col in ["column_left", "column_right"]:
                if not slide.get(col):
                    slide[col] = {}
                if not slide[col].get("title"):
                    slide[col]["title"] = "Option"
                if not slide[col].get("bullets") or len(slide[col].get("bullets", [])) == 0:
                    slide[col]["bullets"] = ["Data point here"]
        elif stype == "timeline":
            if not slide.get("steps") or len(slide.get("steps", [])) == 0:
                slide["steps"] = [{"step": "Step 1", "text": "First phase"},
                                   {"step": "Step 2", "text": "Second phase"}]
        elif stype in ["grid_list", "bento_grid"]:
            if not slide.get("items") or len(slide.get("items", [])) == 0:
                slide["items"] = [{"item_title": "Item", "item_text": "Description",
                                    "title": "Item", "desc": "Description", "size": "small"}]
        elif stype == "metric_dashboard":
            if not slide.get("metrics") or len(slide.get("metrics", [])) == 0:
                slide["metrics"] = [{"label": "Metric", "value": "N/A", "change": "+0%"}]
        elif stype in ["chart_pie", "chart_bar"]:
            if not slide.get("chart_data") or not slide["chart_data"].get("labels"):
                slide["chart_data"] = {"labels": ["Data A", "Data B"], "values": [50, 50]}
        elif stype == "data_table":
            if not slide.get("table_data") or not slide["table_data"].get("headers"):
                slide["table_data"] = {"headers": ["Column 1", "Column 2"],
                                        "rows": [["N/A", "N/A"]]}
    return slides


# ─────────────────────────────────────────────────────────────────────────────
# IMAGE HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def save_base64_image(b64_str):
    if b64_str.startswith("data:image"):
        b64_str = b64_str.split(",")[1]
    img_data = base64.b64decode(b64_str)
    fd, path = tempfile.mkstemp(suffix=".png")
    with os.fdopen(fd, 'wb') as f:
        f.write(img_data)
    return path


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS — each mirrors the corresponding React layout component
# ─────────────────────────────────────────────────────────────────────────────

def build_title_hero_slide(slide, data, theme):
    accent   = theme_color(theme, "accent",    "7C3AED")
    txt_col  = theme_color(theme, "titleColor", "FFFFFF")
    title    = data.get("title", "Presentation Title")
    subtitle = data.get("subtitle", "")

    title_y = SLIDE_H / 2 - 1.0

    # Issue 2 fix: clean accent lines instead of alpha-injected halo bars
    # Top accent line (full width)
    add_rectangle(slide, CONTENT_X, title_y - 0.15, CONTENT_W, 0.04, accent)

    add_text_box(slide, title, 1.0, title_y, CONTENT_W - 0.5, 1.1,
                 40, txt_col, PP_ALIGN.CENTER, bold=True)

    # Bottom accent line
    add_rectangle(slide, CONTENT_X, title_y + 1.18, CONTENT_W, 0.04, accent)

    if subtitle:
        add_text_box(slide, subtitle, 1.5, title_y + 1.35, CONTENT_W - 1.5, 0.7,
                     18, accent, PP_ALIGN.CENTER)


def build_comparison_slide(slide, data, theme):
    """
    Mirrors ComparisonLayout.jsx:
    - Title at top-center
    - Two rounded-rect cards side by side
    - Thin accent-colored top bar on each card
    - "VS" circle badge centered between cards
    - ✓ bullets in left card, → bullets in right card
    """
    accent   = theme_color(theme, "accent",    "7C3AED")
    shape_bg = theme_color(theme, "shapeFill", "1A1A2E")
    border   = theme_color(theme, "bkgd",      "2D3748")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"

    title = data.get("title", "Comparison")
    t_h, b_y, b_h = calc_title_layout(title, base_font_pt=30)
    title_col = theme_color(theme, "titleColor", txt_col)
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, t_h,
                 30, title_col, PP_ALIGN.CENTER, bold=True)

    # Recompute card geometry with dynamic body top
    card_y_base = b_y

    left_d  = data.get("column_left",  {"title": "Option A", "bullets": []})
    right_d = data.get("column_right", {"title": "Option B", "bullets": []})

    # Card dimensions (use dynamic body position from title layout)
    card_y = card_y_base + 0.1
    card_h = SLIDE_H - card_y - MARGIN - 0.1
    card_w = 5.8
    left_x  = CONTENT_X + 0.2
    right_x = SLIDE_W - MARGIN - 0.2 - card_w

    # Left card
    add_rounded_rect(slide, left_x, card_y, card_w, card_h, shape_bg, border)
    # Accent top bar on left card (4pt height ≈ 0.055")
    add_rectangle(slide, left_x, card_y, card_w, 0.06, accent)
    add_text_box(slide, left_d.get("title", "Option A"),
                 left_x + 0.3, card_y + 0.25, card_w - 0.6, 0.55,
                 16, txt_col, PP_ALIGN.CENTER, bold=True)
    add_bullets(slide, left_d.get("bullets", []),
                left_x + 0.3, card_y + 0.9, card_w - 0.6, card_h - 1.1,
                14, sub_col, prefix="✓  ")

    # Right card
    add_rounded_rect(slide, right_x, card_y, card_w, card_h, shape_bg, border)
    add_rectangle(slide, right_x, card_y, card_w, 0.06, accent)
    add_text_box(slide, right_d.get("title", "Option B"),
                 right_x + 0.3, card_y + 0.25, card_w - 0.6, 0.55,
                 16, txt_col, PP_ALIGN.CENTER, bold=True)
    add_bullets(slide, right_d.get("bullets", []),
                right_x + 0.3, card_y + 0.9, card_w - 0.6, card_h - 1.1,
                14, sub_col, prefix="→  ")

    # VS badge — circle in the center gap
    badge_d  = 0.7
    badge_cx = SLIDE_W / 2 - badge_d / 2
    badge_cy = SLIDE_H / 2 - badge_d / 2
    add_oval(slide, badge_cx, badge_cy, badge_d, accent)
    add_text_box(slide, "VS",
                 badge_cx, badge_cy + 0.08, badge_d, badge_d - 0.15,
                 12, "FFFFFF", PP_ALIGN.CENTER, bold=True)


def build_bento_grid_slide(slide, data, theme):
    """
    Mirrors BentoGridLayout.jsx:
    - 3-column grid, items have size: small/large/wide/tall
    - Each cell is a rounded rect with accent top bar
    """
    accent   = theme_color(theme, "accent",    "7C3AED")
    shape_bg = theme_color(theme, "shapeFill", "1A1A2E")
    border   = theme_color(theme, "bkgd",      "2D3748")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"

    title = data.get("title", "Overview")
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, TITLE_H,
                 24, txt_col, PP_ALIGN.CENTER, bold=True)

    items = data.get("items", [])
    if not items:
        return

    # Bug 3 fix: compute rows dynamically from total span area so extra items
    # get a 3rd/4th row instead of being silently dropped.
    cols = 3
    total_span = sum(
        4 if (item.get("size", "small").lower() == "large") else
        3 if (item.get("size", "").lower() == "wide") else
        2 if (item.get("size", "").lower() == "tall") else 1
        for item in items
    )
    rows = max(2, math.ceil(total_span / cols))
    gap      = 0.18
    grid_x   = CONTENT_X + 0.1
    grid_y   = BODY_Y + 0.1
    grid_w   = CONTENT_W - 0.2
    grid_h   = BODY_H - 0.2
    cell_w   = (grid_w - gap * (cols - 1)) / cols
    cell_h   = (grid_h - gap * (rows - 1)) / rows

    # Simple placement: lay items left-to-right, top-to-bottom with size-based spans
    # Track occupied cells with a grid bitmap
    occupied = [[False] * cols for _ in range(rows)]

    def find_free(span_c, span_r):
        for r in range(rows):
            for c in range(cols):
                if c + span_c > cols or r + span_r > rows:
                    continue
                ok = all(not occupied[r + dr][c + dc]
                         for dr in range(span_r) for dc in range(span_c))
                if ok:
                    return r, c
        return None

    def mark_occupied(r, c, span_c, span_r):
        for dr in range(span_r):
            for dc in range(span_c):
                if r + dr < rows and c + dc < cols:
                    occupied[r + dr][c + dc] = True

    for item in items:
        size = (item.get("size") or "small").lower()
        if size == "large":
            span_c, span_r = 2, 2
        elif size == "wide":
            span_c, span_r = 3, 1
        elif size == "tall":
            span_c, span_r = 1, 2
        else:
            span_c, span_r = 1, 1

        pos = find_free(span_c, span_r)
        if pos is None:
            # Fallback: small cell
            span_c, span_r = 1, 1
            pos = find_free(1, 1)
        if pos is None:
            continue

        r, c = pos
        mark_occupied(r, c, span_c, span_r)

        cx = grid_x + c * (cell_w + gap)
        cy = grid_y + r * (cell_h + gap)
        w  = cell_w * span_c + gap * (span_c - 1)
        h  = cell_h * span_r + gap * (span_r - 1)

        # Card
        add_rounded_rect(slide, cx, cy, w, h, shape_bg, border, line_pt=1.0)
        # Accent top stripe
        add_rectangle(slide, cx, cy, w, 0.05, accent)

        title_text = item.get("item_title") or item.get("title", "")
        desc_text  = item.get("item_text")  or item.get("desc",  "")

        # Scale fonts proportionally with cell area — adjusted for extra rows
        area = w * h
        if area > 10:
            fs_title, fs_desc = 22, 14
        elif area > 5:
            fs_title, fs_desc = 18, 13
        elif area > 2:
            fs_title, fs_desc = 13, 11
        else:          # very small cells in high-row grids
            fs_title, fs_desc = 11, 9

        title_box_h = min(0.55, h * 0.2)
        add_text_box(slide, title_text,
                     cx + 0.2, cy + 0.2, w - 0.4, title_box_h,
                     fs_title, txt_col, bold=True)
        add_text_box(slide, desc_text,
                     cx + 0.2, cy + 0.2 + title_box_h + 0.1, w - 0.4, h - 0.2 - title_box_h - 0.2,
                     fs_desc, sub_col)


def build_timeline_slide(slide, data, theme):
    """
    Mirrors TimelineLayout.jsx:
    - Horizontal connector line across mid-slide
    - Even-indexed steps: card ABOVE the line
    - Odd-indexed steps:  card BELOW the line
    - Filled circle dot at each intersection
    """
    accent   = theme_color(theme, "accent",    "7C3AED")
    shape_bg = theme_color(theme, "shapeFill", "1A1A2E")
    border   = theme_color(theme, "bkgd",      "2D3748")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"

    title = data.get("title", "Timeline")
    t_h, b_y, b_h = calc_title_layout(title, base_font_pt=24)
    title_col = theme_color(theme, "titleColor", txt_col)
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, t_h,
                 24, title_col, PP_ALIGN.CENTER, bold=True)

    steps = data.get("steps", [])
    if not steps:
        return
    # Issue 14 fix: cap to 5 steps to avoid card overflow
    steps = steps[:5]

    # Mid-line y position
    line_y = 3.9
    line_x = CONTENT_X + 0.3
    line_w = CONTENT_W - 0.6

    # Draw connector line
    line = add_rectangle(slide, line_x, line_y, line_w, 0.04, sub_col)

    n = len(steps)
    step_w_total = line_w / n
    card_w = min(step_w_total - 0.2, 2.4)
    card_h = 1.5
    dot_d  = 0.28
    # Issue 14 fix: reduce font size for dense timelines to prevent text overflow
    step_font = 10 if n <= 3 else (9 if n == 4 else 8)

    above_card_y = line_y - card_h - 0.35
    below_card_y = line_y + 0.35

    for i, step in enumerate(steps):
        # Horizontal center of this step's dot
        dot_cx = line_x + (i + 0.5) * step_w_total
        card_x = dot_cx - card_w / 2
        dot_left = dot_cx - dot_d / 2
        dot_top  = line_y - dot_d / 2

        # Card position alternates above/below
        card_y = above_card_y if i % 2 == 0 else below_card_y

        # Card
        add_rounded_rect(slide, card_x, card_y, card_w, card_h, shape_bg, border, line_pt=1.0)
        # Step label
        add_text_box(slide, step.get("step", ""),
                     card_x + 0.12, card_y + 0.1, card_w - 0.24, 0.38,
                     step_font, accent, PP_ALIGN.CENTER, bold=True)
        # Step text
        add_text_box(slide, step.get("text", ""),
                     card_x + 0.12, card_y + 0.5, card_w - 0.24, card_h - 0.55,
                     step_font, sub_col, PP_ALIGN.CENTER)

        # Connector vertical stub from dot to card
        if i % 2 == 0:
            stub_y = card_y + card_h
            stub_h = dot_top - stub_y
        else:
            stub_y = dot_top + dot_d
            stub_h = card_y - stub_y
        if stub_h > 0:
            add_rectangle(slide, dot_cx - 0.015, stub_y, 0.03, stub_h, sub_col)

        # Dot
        add_oval(slide, dot_left, dot_top, dot_d, accent)


def build_stat_callout_slide(slide, data, theme):
    """
    Mirrors StatCalloutLayout.jsx:
    - Left half: giant stat number + label (accent colored)
    - Vertical divider line at center
    - Right half: title + bullet list
    """
    accent   = theme_color(theme, "accent",    "7C3AED")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"

    stat  = data.get("stat")  or data.get("huge_text", "100%")
    label = data.get("label") or data.get("subtext",   "")

    # Dynamically size the stat number
    stat_len = len(str(stat))
    if stat_len <= 6:
        stat_pt = 64
    elif stat_len <= 12:
        stat_pt = 48
    else:
        stat_pt = 36

    half_w = CONTENT_W / 2 - 0.4

    # Left — stat number
    stat_y = SLIDE_H / 2 - 1.0
    add_text_box(slide, stat,
                 CONTENT_X, stat_y, half_w, 1.8,
                 stat_pt, accent, PP_ALIGN.CENTER, bold=True)
    if label:
        add_text_box(slide, label.upper(),
                     CONTENT_X, stat_y + 1.9, half_w, 0.6,
                     14, sub_col, PP_ALIGN.CENTER, bold=True)

    # Vertical divider at center
    div_x = SLIDE_W / 2 - 0.01
    add_rectangle(slide, div_x, BODY_Y, 0.03, BODY_H, accent)

    # Right — title + bullets
    right_x = SLIDE_W / 2 + 0.35
    right_w = SLIDE_W - right_x - MARGIN
    title = data.get("title", "")
    add_text_box(slide, title, right_x, BODY_Y + 0.1, right_w, 0.7,
                 18, txt_col, bold=True)

    bullets = data.get("bullets", [])
    if bullets:
        add_bullets(slide, bullets,
                    right_x, BODY_Y + 0.9, right_w, BODY_H - 1.0,
                    13, sub_col, prefix="• ", line_spacing_pt=18)


def build_bento_grid_slide_v2(slide, data, theme):
    """Alias kept for compatibility — calls main bento builder."""
    build_bento_grid_slide(slide, data, theme)


def build_grid_list_slide(slide, data, theme):
    """
    Mirrors GridListLayout.jsx:
    - Title at top
    - 2×2 grid of item cards, each with a numbered badge
    """
    accent   = theme_color(theme, "accent",    "7C3AED")
    shape_bg = theme_color(theme, "shapeFill", "1A1A2E")
    border   = theme_color(theme, "bkgd",      "2D3748")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"

    title = data.get("title", "Overview")
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, TITLE_H,
                 24, txt_col, PP_ALIGN.CENTER, bold=True)

    # Bug 2 fix: Remove [:4] cap — compute rows dynamically so all items render.
    items = data.get("items", [])
    if not items:
        return

    gap   = 0.2
    col_w = (CONTENT_W - gap) / 2
    rows  = math.ceil(len(items) / 2)
    row_h = (BODY_H - gap * (rows - 1)) / rows

    # Scale fonts down when cards become short so text doesn't overflow
    fs_title = 14 if row_h >= 1.4 else (12 if row_h >= 1.0 else 10)
    fs_desc  = 11 if row_h >= 1.4 else (10 if row_h >= 1.0 else 9)

    for i, item in enumerate(items):
        row = i // 2
        col = i % 2
        cx = CONTENT_X + col * (col_w + gap)
        cy = BODY_Y + row * (row_h + gap)

        add_rounded_rect(slide, cx, cy, col_w, row_h, shape_bg, border, line_pt=1.0)
        # Number badge
        add_oval(slide, cx + 0.2, cy + 0.18, 0.38, accent)
        add_text_box(slide, str(i + 1),
                     cx + 0.2, cy + 0.2, 0.38, 0.32,
                     11, "FFFFFF", PP_ALIGN.CENTER, bold=True)
        add_text_box(slide, item.get("item_title") or item.get("title", ""),
                     cx + 0.7, cy + 0.22, col_w - 0.9, 0.45,
                     fs_title, txt_col, bold=True)
        add_text_box(slide, item.get("item_text") or item.get("desc", ""),
                     cx + 0.2, cy + 0.75, col_w - 0.4, row_h - 0.95,
                     fs_desc, sub_col)


def build_three_card_grid_slide(slide, data, theme):
    accent   = theme_color(theme, "accent",    "7C3AED")
    shape_bg = theme_color(theme, "shapeFill", "1A1A2E")
    border   = theme_color(theme, "bkgd",      "2D3748")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"

    title = data.get("title", "Overview")
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, TITLE_H,
                 24, txt_col, PP_ALIGN.CENTER, bold=True)

    # Bug 1 fix: Remove [:3] cap — support 2–6+ cards with dynamic layout.
    cards = data.get("cards", [])
    if not cards:
        return

    n    = len(cards)
    gap  = 0.2
    # For > 4 cards, wrap into 2 rows to keep cards legible
    cols = min(n, 4)
    rows = math.ceil(n / cols)
    card_w = (CONTENT_W - gap * (cols - 1)) / cols
    card_h = (BODY_H - gap * (rows - 1)) / rows - 0.05

    # Scale font sizes down proportionally for wide layouts
    fs_title = max(9,  14 - max(0, n - 3))
    fs_desc  = max(8,  11 - max(0, n - 3))

    for i, card in enumerate(cards):
        row = i // cols
        col = i % cols
        cx = CONTENT_X + col * (card_w + gap)
        cy = BODY_Y + 0.1 + row * (card_h + gap)
        add_rounded_rect(slide, cx, cy, card_w, card_h, shape_bg, border, line_pt=1.0)
        add_rectangle(slide, cx, cy, card_w, 0.05, accent)

        add_text_box(slide, card.get("card_title", ""),
                     cx + 0.2, cy + 0.25, card_w - 0.4, 0.5,
                     fs_title, txt_col, bold=True)
        add_text_box(slide, card.get("card_text", ""),
                     cx + 0.2, cy + 0.8, card_w - 0.4, card_h - 0.9,
                     fs_desc, sub_col)


def build_metric_dashboard_slide(slide, data, theme):
    """
    Mirrors MetricDashboardLayout — row of metric cards:
    label / value / change-badge
    """
    accent   = theme_color(theme, "accent",    "7C3AED")
    shape_bg = theme_color(theme, "shapeFill", "1A1A2E")
    border   = theme_color(theme, "bkgd",      "2D3748")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    green    = "10B981"
    red      = "EF4444"

    title = data.get("title", "Metrics")
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, TITLE_H,
                 24, txt_col, PP_ALIGN.CENTER, bold=True)

    metrics = data.get("metrics", [])
    if not metrics:
        return

    n      = len(metrics)
    gap    = 0.2
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = BODY_H - 0.3
    card_y = BODY_Y + 0.15

    for i, m in enumerate(metrics):
        cx = CONTENT_X + i * (card_w + gap)
        add_rounded_rect(slide, cx, card_y, card_w, card_h, shape_bg, border, line_pt=1.0)
        add_rectangle(slide, cx, card_y, card_w, 0.05, accent)

        # Label
        add_text_box(slide, m.get("label", ""),
                     cx + 0.15, card_y + 0.3, card_w - 0.3, 0.5,
                     12, txt_col, PP_ALIGN.CENTER, bold=True)
        # Value
        add_text_box(slide, m.get("value", ""),
                     cx + 0.1, card_y + 0.9, card_w - 0.2, 1.0,
                     30, accent, PP_ALIGN.CENTER, bold=True)
        # Issue 12 fix: Change badge — red if negative, gray if neutral/zero/N/A, green otherwise
        change = m.get("change", "")
        change_str = str(change).strip()
        neutral_vals = {"+0%", "0%", "0", "n/a", "na", "-", ""}
        if change_str.startswith("-"):
            change_col = red
        elif change_str.lower() in neutral_vals:
            change_col = "6B7280"   # gray-500 for neutral
        else:
            change_col = green
        add_text_box(slide, change,
                     cx + 0.15, card_y + 2.0, card_w - 0.3, 0.45,
                     12, change_col, PP_ALIGN.CENTER, bold=True)


def build_split_card_slide(slide, data, theme, img_path):
    """
    Mirrors TwoColumnLayout / hero_split:
    - Master card container shape
    - Title inside card (left side) or above card
    - Left: bullets / text
    - Right: image or subtitle
    """
    accent   = theme_color(theme, "accent",    "7C3AED")
    shape_bg = theme_color(theme, "shapeFill", "1A1A2E")
    border   = theme_color(theme, "bkgd",      "2D3748")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"

    card_y = BODY_Y + 0.1
    card_h = BODY_H - 0.2
    card_w = CONTENT_W

    # Master card container
    add_rounded_rect(slide, CONTENT_X, card_y, card_w, card_h, shape_bg, border, line_pt=1.0)
    add_rectangle(slide, CONTENT_X, card_y, card_w, 0.06, accent)

    inner_w = (card_w - 0.8) / 2
    left_x = CONTENT_X + 0.4
    right_x = left_x + inner_w + 0.4

    # Title inside left column of card
    title = data.get("title", "")
    add_text_box(slide, title, left_x, card_y + 0.3, inner_w, 0.8,
                 22, txt_col, bold=True)

    bullets = data.get("bullets", [])
    if bullets:
        add_bullets(slide, bullets,
                    left_x, card_y + 1.2, inner_w, card_h - 1.4,
                    13, sub_col, prefix="• ")

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(right_x), Inches(card_y + 0.3),
                                 width=Inches(inner_w), height=Inches(card_h - 0.6))
    else:
        subtitle = data.get("subtitle", "")
        if subtitle:
            add_text_box(slide, subtitle, right_x, card_y + 0.3, inner_w, card_h - 0.6,
                         14, sub_col)



def build_chart_slide(slide, data, theme, chart_type):
    """Pie or Bar chart using pptx native chart engine."""
    txt_col  = theme_color(theme, "textColor", "FFFFFF")

    title = data.get("title", "Chart")
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, TITLE_H,
                 24, txt_col, bold=True)

    chart_data_raw = data.get("chart_data", {})
    labels = chart_data_raw.get("labels", ["A", "B", "C"])
    values = chart_data_raw.get("values", [10, 20, 30])

    cdata = CategoryChartData()
    cdata.categories = labels
    cdata.add_series("Series 1", values)

    ctype = XL_CHART_TYPE.PIE if chart_type == "pie" else XL_CHART_TYPE.COLUMN_CLUSTERED

    chart_left = CONTENT_X + 1.0
    chart = slide.shapes.add_chart(
        ctype,
        Inches(chart_left), Inches(BODY_Y),
        Inches(CONTENT_W - 2.0), Inches(BODY_H),
        cdata).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    try:
        chart.legend.font.color.rgb = hex_to_rgb(txt_col)
    except Exception:
        pass


def build_data_table_slide(slide, data, theme):
    """Data table with styled header row."""
    accent   = theme_color(theme, "accent",    "7C3AED")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"
    shape_bg = theme_color(theme, "shapeFill", "1A1A2E")

    title = data.get("title", "Data Table")
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, TITLE_H,
                 24, txt_col, bold=True)

    tdata   = data.get("table_data", {})
    headers = tdata.get("headers", ["Col 1", "Col 2"])
    rows    = tdata.get("rows", [["A", "B"]])

    num_rows = len(rows) + 1
    num_cols = max(len(headers), 1)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(CONTENT_X), Inches(BODY_Y),
        Inches(CONTENT_W), Inches(BODY_H)).table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(header)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("FFFFFF")
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(accent)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if c < num_cols:
                cell = table.cell(r + 1, c)
                cell.text = str(val)
                cell.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(sub_col)
                cell.fill.solid()
                row_bg = shape_bg if r % 2 == 0 else "111827"
                cell.fill.fore_color.rgb = hex_to_rgb(row_bg)


def build_default_slide(slide, data, theme):
    """
    Standard text slide:
    - Title, optional subtitle, bullet list
    """
    accent   = theme_color(theme, "accent",    "7C3AED")
    txt_col  = theme_color(theme, "textColor", "FFFFFF")
    sub_col  = "A0A8BE"

    title = data.get("title", "")
    # Issue 8 fix: use dynamic title height so long titles don't overflow body
    t_h, b_y, b_h = calc_title_layout(title, base_font_pt=24)
    # Issue 9 fix: use titleColor key for title text
    title_col = theme_color(theme, "titleColor", txt_col)
    add_text_box(slide, title, CONTENT_X, TITLE_Y, CONTENT_W, t_h,
                 24, title_col, bold=True)

    # Accent underline
    add_rectangle(slide, CONTENT_X, TITLE_Y + t_h, CONTENT_W * 0.35, 0.04, accent)

    subtitle    = data.get("subtitle", "")
    current_y   = b_y + 0.05
    bullet_height = b_h - 0.1

    if subtitle:
        add_text_box(slide, subtitle, CONTENT_X, current_y, CONTENT_W, 0.55,
                     14, txt_col, bold=True)
        current_y   += 0.65
        bullet_height -= 0.65

    bullets = data.get("bullets", [])
    if bullets:
        add_bullets(slide, bullets,
                    CONTENT_X, current_y, CONTENT_W, bullet_height,
                    13, sub_col, prefix="• ", line_spacing_pt=18)


# ─────────────────────────────────────────────────────────────────────────────
# BACKGROUND ENGINE
# ─────────────────────────────────────────────────────────────────────────────
def apply_custom_background(slide, bg_data, prs):
    bg_type  = bg_data.get("type")
    bg_value = bg_data.get("value", "")

    if bg_type == "solid":
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = hex_to_rgb(bg_value.lstrip("#"))

    elif bg_type == "image" and bg_value.startswith("data:image"):
        try:
            _, enc = bg_value.split(",", 1)
            img_data = base64.b64decode(enc)
            tf = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tf.write(img_data)
            tf.close()
            pic = slide.shapes.add_picture(
                tf.name, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            os.remove(tf.name)
            pic._element.getparent().insert(0, pic._element)
        except Exception as e:
            print("Failed to set image background:", e)

    elif bg_type == "gradient" and bg_value:
        try:
            colors = re.findall(r"#([A-Fa-f0-9]{6})", bg_value)
            if len(colors) >= 2:
                c1 = hex_to_rgb(colors[0])
                c2 = hex_to_rgb(colors[1])
                w, h = 1280, 720
                img = Image.new("RGB", (w, h))
                draw = ImageDraw.Draw(img)
                for y in range(h):
                    r = int(c1[0] + (c2[0] - c1[0]) * y / h)
                    g = int(c1[1] + (c2[1] - c1[1]) * y / h)
                    b = int(c1[2] + (c2[2] - c1[2]) * y / h)
                    draw.line([(0, y), (w, y)], fill=(r, g, b))
                tf = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                img.save(tf.name)
                tf.close()
                pic = slide.shapes.add_picture(
                    tf.name, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
                os.remove(tf.name)
                pic._element.getparent().insert(0, pic._element)
        except Exception as e:
            print("Failed to set gradient background:", e)

    # Overlay
    if bg_type in ["image", "gradient"] and bg_data.get("overlayOpacity", 0) > 0:
        overlay_color = bg_data.get("overlayColor", "#000000").lstrip("#")
        opacity = bg_data.get("overlayOpacity", 0.5)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        set_shape_fill(rect, overlay_color)
        set_shape_no_line(rect)
        try:
            from pptx.oxml import parse_xml
            alpha_val = int((1.0 - opacity) * 100000)
            alpha_xml = (f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/'
                         f'drawingml/2006/main" val="{alpha_val}"/>')
            srgbClr = rect.fill._xPr.solidFill.srgbClr
            if srgbClr is not None:
                srgbClr.append(parse_xml(alpha_xml))
        except Exception as e:
            print("Transparency XML injection failed:", e)
        rect._element.getparent().insert(1, rect._element)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN EXPORT ORCHESTRATOR
# ─────────────────────────────────────────────────────────────────────────────
def export_presentation(data, output_path, custom_bg=None):
    prs = Presentation()

    # Issue 10 fix: dynamic slide coordinates based on actual canvas size
    slide_size = data.get("slideSize", "LAYOUT_16x9")
    if slide_size == "LAYOUT_4x3":
        prs.slide_width  = Inches(10.0)
        prs.slide_height = Inches(7.5)
        actual_w = 10.0
    elif slide_size == "LAYOUT_16x10":
        prs.slide_width  = Inches(12.0)
        prs.slide_height = Inches(7.5)
        actual_w = 12.0
    else:
        prs.slide_width  = Inches(13.333)
        prs.slide_height = Inches(7.5)
        actual_w = 13.333

    # Patch module-level constants to match the actual canvas
    global SLIDE_W, CONTENT_W, BODY_H
    SLIDE_W   = actual_w
    CONTENT_W = SLIDE_W - MARGIN * 2
    BODY_H    = SLIDE_H - BODY_Y - MARGIN

    blank_layout = prs.slide_layouts[6]

    theme = data.get("theme", {
        "bkgd": "0B0F19", "titleColor": "FFFFFF", "textColor": "FFFFFF",
        "accent": "7C3AED", "shapeFill": "1A1A2E"
    })
    bkgd_color = theme_color(theme, "bkgd", "0B0F19")

    raw_slides   = apply_defensive_fallbacks(data.get("slides", []))
    slides_data  = optimize_presentation(raw_slides)
    temp_images  = []

    try:
        for sdata in slides_data:
            slide = prs.slides.add_slide(blank_layout)
            # Remove any inherited placeholder shapes from the master
            purge_placeholders(slide)

            # ── Background ─────────────────────────────────────────────────────
            # Prefer the CLI --custom_bg arg; fall back to customBackground in JSON
            effective_bg = custom_bg or data.get("customBackground")
            if effective_bg:
                apply_custom_background(slide, effective_bg, prs)
            else:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = hex_to_rgb(bkgd_color)

            # ── Image ──────────────────────────────────────────────────────────
            img_path = None
            if sdata.get("image_base64"):
                img_path = save_base64_image(sdata["image_base64"])
                temp_images.append(img_path)

            # ── Route to builder ───────────────────────────────────────────────
            stype = sdata.get("slide_type", "default")

            if stype == "title_hero":
                build_title_hero_slide(slide, sdata, theme)
            elif stype == "comparison":
                build_comparison_slide(slide, sdata, theme)
            elif stype == "bento_grid":
                build_bento_grid_slide(slide, sdata, theme)
            elif stype == "timeline":
                build_timeline_slide(slide, sdata, theme)
            elif stype == "stat_callout":
                build_stat_callout_slide(slide, sdata, theme)
            elif stype == "grid_list":
                build_grid_list_slide(slide, sdata, theme)
            elif stype == "three_card_grid":
                build_three_card_grid_slide(slide, sdata, theme)
            elif stype == "metric_dashboard":
                build_metric_dashboard_slide(slide, sdata, theme)
            elif stype in ("two_column_image", "hero_split"):
                build_split_card_slide(slide, sdata, theme, img_path)
            elif stype == "chart_pie":
                build_chart_slide(slide, sdata, theme, "pie")
            elif stype == "chart_bar":
                build_chart_slide(slide, sdata, theme, "bar")
            elif stype == "data_table":
                build_data_table_slide(slide, sdata, theme)
            else:
                build_default_slide(slide, sdata, theme)

        prs.save(output_path)
    finally:
        for p in temp_images:
            if os.path.exists(p):
                try:
                    os.remove(p)
                except Exception:
                    pass


# ─────────────────────────────────────────────────────────────────────────────
# CLI ENTRY POINT
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--input",     required=True, help="Input JSON file path")
    parser.add_argument("--output",    required=True, help="Output PPTX file path")
    parser.add_argument("--custom_bg", required=False, help="Custom Background JSON string")
    args = parser.parse_args()

    with open(args.input, "r") as f:
        data = json.load(f)

    custom_bg_obj = None
    if args.custom_bg:
        custom_bg_obj = json.loads(args.custom_bg)

    export_presentation(data, args.output, custom_bg_obj)
    print("SUCCESS")
